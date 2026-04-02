#!/usr/bin/env python3
"""Generate a UofSC Beamer-style PowerPoint comparing the Naive Line Filter
to the Anisotropic Stencil Filter (ASF), focused on computational operations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── UofSC Brand ──────────────────────────────────────────────────────────────
GARNET = RGBColor(0x73, 0x00, 0x0A)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GARNET_LIGHT = RGBColor(0x8C, 0x1D, 0x25)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "papers", "anisotropic-lf-proposal",
)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helpers ──────────────────────────────────────────────────────────────────
def _add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _garnet_bar(slide, top=0, height=Inches(0.9)):
    """Add the garnet header bar (Beamer headline analog)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GARNET
    shape.line.fill.background()
    return shape


def _footer_bar(slide):
    """Thin garnet footer bar with department text."""
    h = Inches(0.35)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_H - h, SLIDE_W, h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GARNET
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Department of Mechanical Engineering  \u2022  University of South Carolina"
    run.font.size = Pt(10)
    run.font.color.rgb = WHITE
    run.font.bold = False


def _header_title(slide, text, subtitle=None):
    """Garnet bar with slide title (Beamer frametitle analog)."""
    bar = _garnet_bar(slide, top=0, height=Inches(0.85))
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(12), Inches(0.65))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(16)
        r2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
        r2.font.name = "Calibri"
    _footer_bar(slide)


def _body_text(slide, left, top, width, height):
    """Return a text frame positioned in the body area."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    return tf


def _add_bullet(tf, text, level=0, size=Pt(20), bold=False, color=DARK_GRAY, spacing_before=Pt(6)):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = spacing_before
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"
    return p


def _first_bullet(tf, text, level=0, size=Pt(20), bold=False, color=DARK_GRAY):
    """Use the existing first (empty) paragraph."""
    p = tf.paragraphs[0]
    p.level = level
    p.space_before = Pt(0)
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"
    return p


def _add_table(slide, rows, cols, left, top, width, height):
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    return tbl


def _set_cell(tbl, r, c, text, bold=False, color=DARK_GRAY, size=Pt(16),
              align=PP_ALIGN.CENTER, fill=None):
    cell = tbl.cell(r, c)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 : Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_add_bg(slide, WHITE)

# Full garnet top half
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(4.2))
bar.fill.solid()
bar.fill.fore_color.rgb = GARNET
bar.line.fill.background()

# Title text
tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.5), Inches(2.5))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "GPU-Accelerated Anisotropic Stencil Filter\nfor Real-Time Edge Detection"
run.font.size = Pt(36)
run.font.color.rgb = WHITE
run.font.bold = True
run.font.name = "Calibri"

# Subtitle
p2 = tf.add_paragraph()
p2.space_before = Pt(16)
p2.alignment = PP_ALIGN.LEFT
r2 = p2.add_run()
r2.text = "Computational Complexity Analysis: Naive LF vs. Fused ASF"
r2.font.size = Pt(22)
r2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
r2.font.name = "Calibri"

# Authors
tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.5), Inches(1.5))
tf2 = tx2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
run = p.add_run()
run.text = "Luke Bagan, Jack Vaught, Yi Wang"
run.font.size = Pt(22)
run.font.color.rgb = DARK_GRAY
run.font.bold = True
run.font.name = "Calibri"

p2 = tf2.add_paragraph()
p2.space_before = Pt(8)
run = p2.add_run()
run.text = "Department of Mechanical Engineering\nUniversity of South Carolina"
run.font.size = Pt(18)
run.font.color.rgb = MID_GRAY
run.font.name = "Calibri"

_footer_bar(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 : Outline
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, WHITE)
_header_title(slide, "Outline")

tf = _body_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5.5))
items = [
    "Background: Polynomial Gradient Estimation",
    "The Naive Line Filter \u2014 Per-Pixel Operation Count",
    "Algebraic Collapse: The Fused Stencil",
    "Deduplication Statistics",
    "Per-Pixel Operation Comparison",
    "GPU Memory Analysis",
    "Wall-Clock Benchmarks (NVIDIA A100)",
    "Summary & Key Takeaways",
]
_first_bullet(tf, items[0], size=Pt(22), color=DARK_GRAY)
for item in items[1:]:
    _add_bullet(tf, item, size=Pt(22), color=DARK_GRAY, spacing_before=Pt(12))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 : Background
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, WHITE)
_header_title(slide, "Background: Polynomial Gradient Estimation")

tf = _body_text(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))
_first_bullet(tf, "Model local intensity with a 2D Taylor expansion of order d", size=Pt(20), color=DARK_GRAY)
_add_bullet(tf, "M = (d+1)(d+2)/2 unknown derivative coefficients", level=1, size=Pt(18), color=MID_GRAY)
_add_bullet(tf, "d = 4  \u2192  M = 15 unknowns", level=1, size=Pt(18), color=MID_GRAY)
_add_bullet(tf, "Gather Np neighbor pixels \u2192 overdetermined system  A\u00b7z = b", size=Pt(20), color=DARK_GRAY, spacing_before=Pt(14))
_add_bullet(tf, "Np = 100 with d = 4  \u2192  6.7\u00d7 overdetermined  \u2192  noise averaging", level=1, size=Pt(18), color=MID_GRAY)
_add_bullet(tf, "Solve via Moore\u2013Penrose pseudoinverse:  \u1e91 = A\u207a b", size=Pt(20), color=DARK_GRAY, spacing_before=Pt(14))
_add_bullet(tf, "Normal derivative = row 1 of A\u207a dotted with pixel intensities", level=1, size=Pt(18), color=MID_GRAY)
_add_bullet(tf, "Sweep Ns orientations, pick the one with max |response|", size=Pt(20), color=DARK_GRAY, spacing_before=Pt(14))
_add_bullet(tf, "Key insight: A\u207a depends only on geometry \u2192 precompute once, reuse per pixel", size=Pt(20), bold=True, color=GARNET, spacing_before=Pt(20))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 : Naive Line Filter Operations
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, WHITE)
_header_title(slide, "Naive Line Filter: Per-Pixel Operations",
              subtitle="m = 7, Np = 100, Ns = 18 orientations, d = 4")

# Left: pseudocode box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(1.15), Inches(6.0), Inches(5.1))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_GRAY
box.line.color.rgb = MID_GRAY
box.line.width = Pt(1)

tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(5.5), Inches(4.8))
tf = tx.text_frame
tf.word_wrap = True
code_lines = [
    "for each orientation \u03b8  (Ns = 18):",
    "  for j = \u2212m to m      (L = 2m+1 = 15):",
    "    for i = 1 to Np    (Np = 100):",
    "      READ pixel at virtual position j,",
    "             neighbor i",
    "      MULTIPLY by pseudoinverse weight",
    "    SUM Np weighted values    (99 adds)",
    "  MULTIPLY by Gaussian w_j   (15 muls)",
    "  SUM over line positions     (14 adds)",
    "  |R_\u03b8| = absolute value",
    "max over all \u03b8",
]
p = tf.paragraphs[0]
run = p.add_run()
run.text = code_lines[0]
run.font.size = Pt(16)
run.font.name = "Consolas"
run.font.color.rgb = DARK_GRAY
for line in code_lines[1:]:
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run()
    run.text = line
    run.font.size = Pt(16)
    run.font.name = "Consolas"
    run.font.color.rgb = DARK_GRAY

# Right: operation counts
tf2 = _body_text(slide, Inches(6.8), Inches(1.15), Inches(6.0), Inches(5.5))
_first_bullet(tf2, "Per orientation:", size=Pt(20), bold=True, color=GARNET)
_add_bullet(tf2, "Memory reads:  15 \u00d7 100 = 1,500", size=Pt(19), color=DARK_GRAY)
_add_bullet(tf2, "Multiplies:  15 \u00d7 100 + 15 = 1,515", size=Pt(19), color=DARK_GRAY)
_add_bullet(tf2, "Additions:  15 \u00d7 99 + 14 = 1,499", size=Pt(19), color=DARK_GRAY)

_add_bullet(tf2, "Per pixel (18 orientations):", size=Pt(20), bold=True, color=GARNET, spacing_before=Pt(18))
_add_bullet(tf2, "Memory reads:  27,000", size=Pt(19), color=DARK_GRAY)
_add_bullet(tf2, "Multiplies:    27,270", size=Pt(19), color=DARK_GRAY)
_add_bullet(tf2, "Additions:     26,982", size=Pt(19), color=DARK_GRAY)
_add_bullet(tf2, "Total ops:     \u224881,252", size=Pt(19), bold=True, color=DARK_GRAY)

_add_bullet(tf2, "Complexity:  O(Ns \u00d7 L \u00d7 Np)", size=Pt(20), bold=True, color=GARNET, spacing_before=Pt(18))
_add_bullet(tf2, "= O(Ns \u00d7 (2m+1) \u00d7 Np)", level=1, size=Pt(18), color=MID_GRAY)

_footer_bar(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 : Algebraic Collapse
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, WHITE)
_header_title(slide, "Algebraic Collapse: The Fused Stencil")

tf = _body_text(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(2.0))
_first_bullet(tf, "The line-extended response is a linear operation on pixel intensities:", size=Pt(20), color=DARK_GRAY)

# Equation box
eq_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(1.5), Inches(2.2), Inches(10.0), Inches(1.0))
eq_box.fill.solid()
eq_box.fill.fore_color.rgb = LIGHT_GRAY
eq_box.line.color.rgb = MID_GRAY

eq_tx = slide.shapes.add_textbox(Inches(1.8), Inches(2.3), Inches(9.5), Inches(0.8))
eq_tf = eq_tx.text_frame
eq_tf.word_wrap = True
p = eq_tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "R\u2096 = \u03a3 w\u2c7c \u00b7 p\u1d62\u207c\u1d4f\u207e \u00b7 f(pixel)   \u2192   R\u2096 = \u03b1\u2096\u1d40 \u00b7 g\u2096"
run.font.size = Pt(26)
run.font.name = "Cambria Math"
run.font.color.rgb = GARNET
run.font.bold = True

tf2 = _body_text(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(3.5))
_first_bullet(tf2, "Many (j, i) pairs map to the same pixel after rounding to integer coords", size=Pt(20), color=DARK_GRAY)
_add_bullet(tf2, "Sum their weights \u2192 single precomputed \u03b1 per unique position", size=Pt(20), color=DARK_GRAY, spacing_before=Pt(10))
_add_bullet(tf2, "The entire pipeline collapses into one gather + dot product:", size=Pt(20), color=DARK_GRAY, spacing_before=Pt(10))

_add_bullet(tf2, "Rotate \u2192 Build A \u2192 A\u207a \u2192 Gather \u2192 MatMul \u2192 Weight \u2192 Sum", level=1, size=Pt(18), color=ACCENT_RED, spacing_before=Pt(8))
_add_bullet(tf2, "\u2193", level=1, size=Pt(22), bold=True, color=DARK_GRAY, spacing_before=Pt(2))
_add_bullet(tf2, "Single Gather  +  Dot Product  (precomputed \u03b1\u2096)", level=1, size=Pt(18), color=ACCENT_GREEN, spacing_before=Pt(2))

_add_bullet(tf2, "All pseudoinverse, rotation, and Gaussian weighting absorbed into \u03b1\u2096 at precompute time", size=Pt(19), bold=True, color=GARNET, spacing_before=Pt(16))

_footer_bar(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 : Deduplication Statistics
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, WHITE)
_header_title(slide, "Stencil Deduplication Statistics",
              subtitle="Np = 100 neighbors, d = 4, averaged across all orientations")

tbl = _add_table(slide, 5, 5, Inches(1.5), Inches(1.3), Inches(10.0), Inches(2.5))

headers = ["Half-width m", "Raw stencil size\n(2m+1) \u00d7 Np", "Unique positions\nN\u2032", "Reduction", "Effective\nspeedup factor"]
for c, h in enumerate(headers):
    _set_cell(tbl, 0, c, h, bold=True, color=WHITE, size=Pt(15), fill=GARNET)

data = [
    ("1",  "300",   "128", "57%",  "2.3\u00d7"),
    ("2",  "500",   "152", "70%",  "3.3\u00d7"),
    ("7",  "1,500", "264", "82%",  "5.7\u00d7"),
    ("14", "2,900", "431", "85%",  "6.7\u00d7"),
]
for r, row_data in enumerate(data):
    fill_color = LIGHT_GRAY if r % 2 == 0 else WHITE
    for c, val in enumerate(row_data):
        _set_cell(tbl, r + 1, c, val, size=Pt(17), fill=fill_color)

# Annotation
tf = _body_text(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5))
_first_bullet(tf, "Neighboring virtual positions share most of their circular neighborhoods", size=Pt(20), color=DARK_GRAY)
_add_bullet(tf, "At m = 7: only 264 unique pixel reads instead of 1,500 \u2192 82% fewer memory accesses", size=Pt(20), color=DARK_GRAY)
_add_bullet(tf, "At m = 14: 85% reduction \u2014 deduplication improves as m grows", size=Pt(20), color=DARK_GRAY)
_add_bullet(tf, "This is the primary source of computational savings", size=Pt(20), bold=True, color=GARNET, spacing_before=Pt(14))

_footer_bar(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 : ASF Per-Pixel Operations
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, WHITE)
_header_title(slide, "Fused ASF: Per-Pixel Operations",
              subtitle="m = 7, Np = 100, Ns = 18 orientations, d = 4")

# Left: pseudocode
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(1.15), Inches(6.0), Inches(3.2))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_GRAY
box.line.color.rgb = MID_GRAY

tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(5.5), Inches(3.0))
tf = tx.text_frame
tf.word_wrap = True
code_lines2 = [
    "for each orientation \u03b8  (Ns = 18):",
    "  for i = 1 to N\u2032     (N\u2032 \u2248 264):",
    "    READ pixel at stencil offset i",
    "    MULTIPLY by precomputed \u03b1[i]",
    "  SUM N\u2032 weighted values  (263 adds)",
    "  |R_\u03b8| = absolute value",
    "max over all \u03b8",
]
p = tf.paragraphs[0]
run = p.add_run()
run.text = code_lines2[0]
run.font.size = Pt(17)
run.font.name = "Consolas"
run.font.color.rgb = DARK_GRAY
for line in code_lines2[1:]:
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = line
    run.font.size = Pt(17)
    run.font.name = "Consolas"
    run.font.color.rgb = DARK_GRAY

# Right: counts
tf2 = _body_text(slide, Inches(6.8), Inches(1.15), Inches(6.0), Inches(3.5))
_first_bullet(tf2, "Per orientation:", size=Pt(20), bold=True, color=GARNET)
_add_bullet(tf2, "Memory reads:  264", size=Pt(19), color=DARK_GRAY)
_add_bullet(tf2, "Multiplies:    264", size=Pt(19), color=DARK_GRAY)
_add_bullet(tf2, "Additions:     263", size=Pt(19), color=DARK_GRAY)

_add_bullet(tf2, "Per pixel (18 orientations):", size=Pt(20), bold=True, color=GARNET, spacing_before=Pt(18))
_add_bullet(tf2, "Memory reads:  4,752", size=Pt(19), color=DARK_GRAY)
_add_bullet(tf2, "Multiplies:    4,752", size=Pt(19), color=DARK_GRAY)
_add_bullet(tf2, "Additions:     4,734", size=Pt(19), color=DARK_GRAY)
_add_bullet(tf2, "Total ops:     \u224814,238", size=Pt(19), bold=True, color=DARK_GRAY)

# Bottom: complexity
_add_bullet(tf2, "Complexity:  O(Ns \u00d7 N\u2032)", size=Pt(20), bold=True, color=GARNET, spacing_before=Pt(18))

# No runtime matrix ops callout
callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.9))
callout.fill.solid()
callout.fill.fore_color.rgb = RGBColor(0xF9, 0xEB, 0xEA)
callout.line.color.rgb = GARNET

tx3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.8), Inches(0.7))
tf3 = tx3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Zero runtime matrix operations \u2014 all pseudoinverse, rotation, and Gaussian weights baked into \u03b1\u2096"
run.font.size = Pt(20)
run.font.bold = True
run.font.color.rgb = GARNET
run.font.name = "Calibri"

_footer_bar(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 : Side-by-Side Comparison Table
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, WHITE)
_header_title(slide, "Operation Count Comparison",
              subtitle="Per pixel, m = 7, Np = 100, Ns = 18, d = 4")

tbl = _add_table(slide, 6, 4, Inches(1.2), Inches(1.2), Inches(10.5), Inches(3.5))

headers2 = ["Metric", "Naive LF", "Fused ASF", "Reduction"]
for c, h in enumerate(headers2):
    _set_cell(tbl, 0, c, h, bold=True, color=WHITE, size=Pt(17), fill=GARNET)

rows2 = [
    ("Memory reads / orientation",   "1,500",  "264",    "5.7\u00d7"),
    ("Multiplies / orientation",     "1,515",  "264",    "5.7\u00d7"),
    ("Additions / orientation",      "1,499",  "263",    "5.7\u00d7"),
    ("Total ops / pixel (18 orient)","81,252", "14,238", "5.7\u00d7"),
    ("Runtime matrix ops / pixel",   "18 matmuls", "0", "\u221e"),
]
for r, (label, naive, asf, reduction) in enumerate(rows2):
    fill_color = LIGHT_GRAY if r % 2 == 0 else WHITE
    _set_cell(tbl, r+1, 0, label, size=Pt(16), fill=fill_color, align=PP_ALIGN.LEFT)
    _set_cell(tbl, r+1, 1, naive, size=Pt(16), fill=fill_color, color=ACCENT_RED)
    _set_cell(tbl, r+1, 2, asf, size=Pt(16), fill=fill_color, color=ACCENT_GREEN)
    _set_cell(tbl, r+1, 3, reduction, size=Pt(16), fill=fill_color, bold=True)

# Scaling note
tf = _body_text(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.5))
_first_bullet(tf, "5.7\u00d7 arithmetic reduction from deduplication alone (m = 7)", size=Pt(20), color=DARK_GRAY)
_add_bullet(tf, "Additional gains on GPU: eliminating runtime pseudoinverse \u2192 18.4\u00d7 wall-clock speedup", size=Pt(20), color=DARK_GRAY)
_add_bullet(tf, "At m = 14: 6.7\u00d7 arithmetic reduction \u2192 24\u00d7 wall-clock speedup", size=Pt(20), bold=True, color=GARNET, spacing_before=Pt(10))

_footer_bar(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 : GPU Memory Analysis
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, WHITE)
_header_title(slide, "GPU Memory Analysis",
              subtitle="BIPED v1 (1280\u00d7720), m = 7, Np = 100, Ns = 18")

# Left column: Naive
box_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(1.2), Inches(5.8), Inches(3.5))
box_l.fill.solid()
box_l.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
box_l.line.color.rgb = ACCENT_RED

tx_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(5.3), Inches(3.3))
tf_l = tx_l.text_frame
tf_l.word_wrap = True
_first_bullet(tf_l, "Naive Line Filter", size=Pt(22), bold=True, color=ACCENT_RED)
_add_bullet(tf_l, "Intermediate tensor per batch:", size=Pt(18), color=DARK_GRAY, spacing_before=Pt(12))
_add_bullet(tf_l, "B \u00d7 L \u00d7 Np \u00d7 24 bytes", level=1, size=Pt(17), color=MID_GRAY)
_add_bullet(tf_l, "Grows with m and Np", size=Pt(18), color=DARK_GRAY)
_add_bullet(tf_l, "Peak VRAM: 6,223 MB", size=Pt(20), bold=True, color=ACCENT_RED, spacing_before=Pt(14))
_add_bullet(tf_l, "Cannot fit large m on\nconsumer GPUs", size=Pt(18), color=DARK_GRAY, spacing_before=Pt(8))

# Right column: ASF
box_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(6.8), Inches(1.2), Inches(5.8), Inches(3.5))
box_r.fill.solid()
box_r.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
box_r.line.color.rgb = ACCENT_GREEN

tx_r = slide.shapes.add_textbox(Inches(7.1), Inches(1.35), Inches(5.3), Inches(3.3))
tf_r = tx_r.text_frame
tf_r.word_wrap = True
_first_bullet(tf_r, "Fused ASF (Triton)", size=Pt(22), bold=True, color=ACCENT_GREEN)
_add_bullet(tf_r, "No intermediate tensors:", size=Pt(18), color=DARK_GRAY, spacing_before=Pt(12))
_add_bullet(tf_r, "Image + output + stencil weights", level=1, size=Pt(17), color=MID_GRAY)
_add_bullet(tf_r, "Constant regardless of m", size=Pt(18), color=DARK_GRAY)
_add_bullet(tf_r, "Peak VRAM: 20 MB", size=Pt(20), bold=True, color=ACCENT_GREEN, spacing_before=Pt(14))
_add_bullet(tf_r, "Fits on any GPU,\nincluding embedded", size=Pt(18), color=DARK_GRAY, spacing_before=Pt(8))

# Bottom: 311x callout
callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(2.0), Inches(5.0), Inches(9.0), Inches(1.0))
callout.fill.solid()
callout.fill.fore_color.rgb = GARNET
callout.line.fill.background()

tx_c = slide.shapes.add_textbox(Inches(2.2), Inches(5.1), Inches(8.6), Inches(0.8))
tf_c = tx_c.text_frame
tf_c.word_wrap = True
p = tf_c.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "311\u00d7 VRAM reduction  \u2014  6,223 MB \u2192 20 MB"
run.font.size = Pt(28)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = "Calibri"

_footer_bar(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 : Wall-Clock Benchmarks
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, WHITE)
_header_title(slide, "Wall-Clock Benchmarks",
              subtitle="BIPED v1 (1280\u00d7720), Np = 100, Ns = 18, d = 4, NVIDIA A100")

tbl = _add_table(slide, 5, 5, Inches(0.8), Inches(1.2), Inches(11.5), Inches(2.8))

headers3 = ["Half-width m", "Naive LF (s)", "Conv2d (s)", "Triton ASF (s)", "Speedup"]
for c, h in enumerate(headers3):
    _set_cell(tbl, 0, c, h, bold=True, color=WHITE, size=Pt(16), fill=GARNET)

bench_data = [
    ("1",  "0.58", "0.089", "0.072", "8.0\u00d7"),
    ("2",  "1.91", "0.187", "0.140", "13.7\u00d7"),
    ("7",  "2.43", "0.176", "0.132", "18.4\u00d7"),
    ("14", "8.88", "0.591", "0.366", "24.3\u00d7"),
]
for r, row_data in enumerate(bench_data):
    fill_color = LIGHT_GRAY if r % 2 == 0 else WHITE
    for c, val in enumerate(row_data):
        clr = DARK_GRAY
        if c == 1:
            clr = ACCENT_RED
        elif c == 3:
            clr = ACCENT_GREEN
        elif c == 4:
            clr = GARNET
        _set_cell(tbl, r+1, c, val, size=Pt(17), fill=fill_color,
                  color=clr, bold=(c == 4))

# Key observations
tf = _body_text(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.3))
_first_bullet(tf, "Speedup scales with m because deduplication ratio improves", size=Pt(20), color=DARK_GRAY)
_add_bullet(tf, "Additional speedup beyond dedup ratio comes from eliminating runtime matrix ops", size=Pt(20), color=DARK_GRAY)
_add_bullet(tf, "Triton kernel processes 1280\u00d7720 in 132 ms (m = 7) or 366 ms (m = 14)", size=Pt(20), color=DARK_GRAY)
_add_bullet(tf, "24\u00d7 faster at m = 14, with constant 20 MB VRAM", size=Pt(20), bold=True, color=GARNET, spacing_before=Pt(10))

_footer_bar(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 : Scaling Breakdown
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, WHITE)
_header_title(slide, "Where Do the Savings Come From?")

# Three source boxes
sources = [
    ("1. Deduplication", ACCENT_RED, [
        "Overlapping neighborhoods share pixels",
        "82% reduction at m = 7 (1,500 \u2192 264 reads)",
        "Savings grow with m: 57% (m=1) \u2192 85% (m=14)",
        "\u2192  5.7\u00d7 fewer memory accesses",
    ]),
    ("2. Precomputed Weights", GARNET, [
        "Pseudoinverse, rotation, Gaussian weights",
        "  all fused into single \u03b1\u2096 vector",
        "Zero runtime matrix operations",
        "\u2192  Eliminates O(Np \u00d7 M) matmul per orientation",
    ]),
    ("3. GPU Kernel Design", ACCENT_GREEN, [
        "Single gather + dot product per orientation",
        "Vectorized across 128-pixel blocks",
        "Coalesced memory reads along image rows",
        "\u2192  Constant 20 MB VRAM regardless of m",
    ]),
]

for idx, (title, color, bullets) in enumerate(sources):
    left = Inches(0.4 + idx * 4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, Inches(1.15), Inches(3.9), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        left, Inches(1.15), Inches(3.9), Inches(0.55))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = color
    title_bar.line.fill.background()
    ttx = slide.shapes.add_textbox(left + Inches(0.15), Inches(1.2), Inches(3.6), Inches(0.45))
    ttf = ttx.text_frame
    p = ttf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    btx = slide.shapes.add_textbox(left + Inches(0.15), Inches(1.85), Inches(3.6), Inches(3.9))
    btf = btx.text_frame
    btf.word_wrap = True
    _first_bullet(btf, bullets[0], size=Pt(15), color=DARK_GRAY)
    for b in bullets[1:]:
        _add_bullet(btf, b, size=Pt(15), color=DARK_GRAY, spacing_before=Pt(8))

_footer_bar(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 : Summary
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(slide, WHITE)
_header_title(slide, "Summary & Key Takeaways")

tf = _body_text(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))

_first_bullet(tf, "The naive line filter is memory-bound and redundant", size=Pt(22), bold=True, color=GARNET)
_add_bullet(tf, "O(Ns \u00d7 (2m+1) \u00d7 Np) operations per pixel, with most pixel reads duplicated", level=1, size=Pt(19), color=DARK_GRAY)
_add_bullet(tf, "~81,000 total ops/pixel at typical settings (m=7, Np=100, Ns=18)", level=1, size=Pt(19), color=DARK_GRAY)

_add_bullet(tf, "The fused stencil eliminates all redundancy", size=Pt(22), bold=True, color=GARNET, spacing_before=Pt(20))
_add_bullet(tf, "O(Ns \u00d7 N\u2032) operations with N\u2032 \u2248 0.18 \u00d7 (2m+1) \u00d7 Np", level=1, size=Pt(19), color=DARK_GRAY)
_add_bullet(tf, "~14,000 total ops/pixel \u2014 5.7\u00d7 arithmetic reduction", level=1, size=Pt(19), color=DARK_GRAY)

_add_bullet(tf, "GPU results exceed the theoretical arithmetic speedup", size=Pt(22), bold=True, color=GARNET, spacing_before=Pt(20))
_add_bullet(tf, "18.4\u00d7 wall-clock speedup at m = 7, 24\u00d7 at m = 14", level=1, size=Pt(19), color=DARK_GRAY)
_add_bullet(tf, "311\u00d7 VRAM reduction (6.2 GB \u2192 20 MB) enables deployment on any GPU", level=1, size=Pt(19), color=DARK_GRAY)
_add_bullet(tf, "Extra gains from eliminating runtime matrix ops + GPU kernel optimization", level=1, size=Pt(19), color=DARK_GRAY)

_footer_bar(slide)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out_path = os.path.join(OUT_DIR, "ASF_Computational_Comparison.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
